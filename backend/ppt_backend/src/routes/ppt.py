from flask import Blueprint, jsonify, request, send_file
from flask_cors import cross_origin
import os
import json
import uuid
import base64
from io import BytesIO
from PIL import Image
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import openai

ppt_bp = Blueprint('ppt', __name__)

# Global variable to store session data
sessions = {}

def call_llm(prompt, api_key, endpoint, model="gpt-3.5-turbo"):
    """Call LLM API with user's configuration"""
    try:
        if "openai" in endpoint.lower() or "api.openai.com" in endpoint:
            # OpenAI API
            client = openai.OpenAI(api_key=api_key, base_url=endpoint if not endpoint.endswith("openai.com") else None)
            response = client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=2000
            )
            return response.choices[0].message.content
        else:
            # Generic API endpoint
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            data = {
                "model": model,
                "messages": [{"role": "user", "content": prompt}],
                "max_tokens": 2000
            }
            response = requests.post(endpoint, headers=headers, json=data)
            if response.status_code == 200:
                return response.json()["choices"][0]["message"]["content"]
            else:
                return f"Error: {response.status_code} - {response.text}"
    except Exception as e:
        return f"Error calling LLM: {str(e)}"

def generate_ppt_structure(user_input, api_key, endpoint, model):
    """Generate PPT structure using LLM"""
    prompt = f"""
    Based on the following input, create a detailed PowerPoint presentation structure:
    
    Input: {user_input}
    
    Please provide a JSON response with the following structure:
    {{
        "title": "Main presentation title",
        "theme": {{
            "primary_color": "#hex_color",
            "secondary_color": "#hex_color",
            "background_color": "#hex_color",
            "text_color": "#hex_color"
        }},
        "slides": [
            {{
                "title": "Slide title",
                "content": "Main content for the slide",
                "bullet_points": ["Point 1", "Point 2", "Point 3"],
                "image_suggestion": "Description of relevant image"
            }}
        ]
    }}
    
    Create 5-8 slides maximum. Make the content engaging and professional.
    """
    
    response = call_llm(prompt, api_key, endpoint, model)
    try:
        # Extract JSON from response if it's wrapped in markdown
        if "```json" in response:
            response = response.split("```json")[1].split("```")[0]
        elif "```" in response:
            response = response.split("```")[1].split("```")[0]
        
        return json.loads(response.strip())
    except json.JSONDecodeError:
        # Fallback structure if JSON parsing fails
        return {
            "title": "Generated Presentation",
            "theme": {
                "primary_color": "#2E86AB",
                "secondary_color": "#A23B72",
                "background_color": "#F18F01",
                "text_color": "#C73E1D"
            },
            "slides": [
                {
                    "title": "Introduction",
                    "content": "Welcome to our presentation",
                    "bullet_points": ["Overview", "Objectives", "Agenda"],
                    "image_suggestion": "Professional business image"
                }
            ]
        }

def create_ppt_from_structure(structure, session_id):
    """Create PowerPoint file from structure"""
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = structure["title"]
    subtitle.text = "Generated Presentation"
    
    # Content slides
    for slide_data in structure["slides"]:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = slide_data["title"]
        
        # Add content
        tf = content.text_frame
        tf.text = slide_data["content"]
        
        # Add bullet points
        if "bullet_points" in slide_data:
            for point in slide_data["bullet_points"]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 1
    
    # Save PPT
    output_dir = f"/tmp/ppt_sessions/{session_id}"
    os.makedirs(output_dir, exist_ok=True)
    ppt_path = os.path.join(output_dir, "presentation.pptx")
    prs.save(ppt_path)
    
    return ppt_path

def generate_thumbnails(ppt_path, session_id):
    """Generate thumbnail images from PPT slides"""
    try:
        from pptx import Presentation
        prs = Presentation(ppt_path)
        
        thumbnails = []
        output_dir = f"/tmp/ppt_sessions/{session_id}/thumbnails"
        os.makedirs(output_dir, exist_ok=True)
        
        for i, slide in enumerate(prs.slides):
            # Create a simple thumbnail representation
            # Note: This is a simplified version. In production, you might want to use
            # libraries like python-pptx-interface or convert to images using other tools
            thumbnail_path = os.path.join(output_dir, f"slide_{i+1}.png")
            
            # Create a simple placeholder image for now
            img = Image.new('RGB', (320, 240), color='white')
            img.save(thumbnail_path)
            
            thumbnails.append({
                "slide_number": i + 1,
                "thumbnail_path": thumbnail_path,
                "thumbnail_url": f"/api/ppt/thumbnail/{session_id}/{i+1}"
            })
        
        return thumbnails
    except Exception as e:
        print(f"Error generating thumbnails: {e}")
        return []

@ppt_bp.route('/generate', methods=['POST'])
@cross_origin()
def generate_ppt():
    """Generate PPT from user input"""
    try:
        data = request.json
        user_input = data.get('input', '')
        api_key = data.get('api_key', '')
        endpoint = data.get('endpoint', 'https://api.openai.com/v1/chat/completions')
        model = data.get('model', 'gpt-3.5-turbo')
        session_id = data.get('session_id', str(uuid.uuid4()))
        
        if not user_input or not api_key:
            return jsonify({"error": "Input and API key are required"}), 400
        
        # Generate PPT structure using LLM
        structure = generate_ppt_structure(user_input, api_key, endpoint, model)
        
        # Create PPT file
        ppt_path = create_ppt_from_structure(structure, session_id)
        
        # Generate thumbnails
        thumbnails = generate_thumbnails(ppt_path, session_id)
        
        # Store session data
        sessions[session_id] = {
            "structure": structure,
            "ppt_path": ppt_path,
            "thumbnails": thumbnails,
            "api_key": api_key,
            "endpoint": endpoint,
            "model": model,
            "original_input": user_input
        }
        
        return jsonify({
            "session_id": session_id,
            "structure": structure,
            "thumbnails": thumbnails,
            "download_url": f"/api/ppt/download/{session_id}"
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@ppt_bp.route('/refine', methods=['POST'])
@cross_origin()
def refine_ppt():
    """Refine existing PPT with additional prompts"""
    try:
        data = request.json
        session_id = data.get('session_id', '')
        refinement_prompt = data.get('prompt', '')
        
        if not session_id or session_id not in sessions:
            return jsonify({"error": "Invalid session ID"}), 400
        
        if not refinement_prompt:
            return jsonify({"error": "Refinement prompt is required"}), 400
        
        session_data = sessions[session_id]
        
        # Create refinement prompt with context
        context_prompt = f"""
        Original input: {session_data['original_input']}
        Current presentation structure: {json.dumps(session_data['structure'], indent=2)}
        
        Refinement request: {refinement_prompt}
        
        Please provide an updated JSON structure for the presentation incorporating the refinement request.
        Keep the same JSON format as before.
        """
        
        # Get refined structure
        refined_structure = generate_ppt_structure(
            context_prompt, 
            session_data['api_key'], 
            session_data['endpoint'], 
            session_data['model']
        )
        
        # Create new PPT file
        ppt_path = create_ppt_from_structure(refined_structure, session_id)
        
        # Generate new thumbnails
        thumbnails = generate_thumbnails(ppt_path, session_id)
        
        # Update session data
        sessions[session_id].update({
            "structure": refined_structure,
            "ppt_path": ppt_path,
            "thumbnails": thumbnails
        })
        
        return jsonify({
            "session_id": session_id,
            "structure": refined_structure,
            "thumbnails": thumbnails,
            "download_url": f"/api/ppt/download/{session_id}"
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@ppt_bp.route('/download/<session_id>', methods=['GET'])
@cross_origin()
def download_ppt(session_id):
    """Download generated PPT file"""
    try:
        if session_id not in sessions:
            return jsonify({"error": "Session not found"}), 404
        
        ppt_path = sessions[session_id]["ppt_path"]
        if not os.path.exists(ppt_path):
            return jsonify({"error": "PPT file not found"}), 404
        
        return send_file(
            ppt_path,
            as_attachment=True,
            download_name="generated_presentation.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@ppt_bp.route('/thumbnail/<session_id>/<int:slide_number>', methods=['GET'])
@cross_origin()
def get_thumbnail(session_id, slide_number):
    """Get thumbnail image for a specific slide"""
    try:
        if session_id not in sessions:
            return jsonify({"error": "Session not found"}), 404
        
        thumbnail_path = f"/tmp/ppt_sessions/{session_id}/thumbnails/slide_{slide_number}.png"
        if not os.path.exists(thumbnail_path):
            return jsonify({"error": "Thumbnail not found"}), 404
        
        return send_file(thumbnail_path, mimetype="image/png")
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@ppt_bp.route('/upload', methods=['POST'])
@cross_origin()
def upload_file():
    """Handle file uploads (images or existing PPTs)"""
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        session_id = str(uuid.uuid4())
        upload_dir = f"/tmp/ppt_sessions/{session_id}/uploads"
        os.makedirs(upload_dir, exist_ok=True)
        
        file_path = os.path.join(upload_dir, file.filename)
        file.save(file_path)
        
        # Process based on file type
        file_ext = file.filename.lower().split('.')[-1]
        
        if file_ext in ['jpg', 'jpeg', 'png', 'gif']:
            # Image file - analyze and suggest PPT content
            return jsonify({
                "session_id": session_id,
                "file_type": "image",
                "file_path": file_path,
                "message": "Image uploaded successfully. Please provide your LLM configuration to analyze the image."
            })
        elif file_ext in ['ppt', 'pptx']:
            # PowerPoint file - extract content for refinement
            return jsonify({
                "session_id": session_id,
                "file_type": "presentation",
                "file_path": file_path,
                "message": "Presentation uploaded successfully. Please provide your LLM configuration to analyze and refine the presentation."
            })
        else:
            return jsonify({"error": "Unsupported file type"}), 400
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@ppt_bp.route('/analyze-upload', methods=['POST'])
@cross_origin()
def analyze_upload():
    """Analyze uploaded file and generate PPT suggestions"""
    try:
        data = request.json
        session_id = data.get('session_id', '')
        api_key = data.get('api_key', '')
        endpoint = data.get('endpoint', 'https://api.openai.com/v1/chat/completions')
        model = data.get('model', 'gpt-3.5-turbo')
        
        if not session_id or not api_key:
            return jsonify({"error": "Session ID and API key are required"}), 400
        
        upload_dir = f"/tmp/ppt_sessions/{session_id}/uploads"
        if not os.path.exists(upload_dir):
            return jsonify({"error": "No uploaded files found"}), 404
        
        # Find uploaded file
        uploaded_files = os.listdir(upload_dir)
        if not uploaded_files:
            return jsonify({"error": "No uploaded files found"}), 404
        
        file_path = os.path.join(upload_dir, uploaded_files[0])
        file_ext = uploaded_files[0].lower().split('.')[-1]
        
        if file_ext in ['jpg', 'jpeg', 'png', 'gif']:
            # Analyze image and suggest PPT content
            prompt = """
            Based on an uploaded image, create a PowerPoint presentation structure.
            Since I cannot see the image directly, please create a general business presentation structure
            that could work with various types of images. Include slides for:
            1. Title slide
            2. Overview/Introduction
            3. Key Points (3-4 slides)
            4. Conclusion
            
            Provide the response in the same JSON format as before.
            """
        else:
            # Analyze existing PPT (simplified - in production you'd extract actual content)
            prompt = """
            Create an improved version of an existing PowerPoint presentation.
            Since I cannot access the original content directly, please create a modern,
            professional presentation structure with:
            1. Enhanced title slide
            2. Executive summary
            3. Main content sections (4-5 slides)
            4. Key takeaways
            5. Next steps/conclusion
            
            Provide the response in the same JSON format as before.
            """
        
        # Generate structure
        structure = generate_ppt_structure(prompt, api_key, endpoint, model)
        
        # Create PPT file
        ppt_path = create_ppt_from_structure(structure, session_id)
        
        # Generate thumbnails
        thumbnails = generate_thumbnails(ppt_path, session_id)
        
        # Store session data
        sessions[session_id] = {
            "structure": structure,
            "ppt_path": ppt_path,
            "thumbnails": thumbnails,
            "api_key": api_key,
            "endpoint": endpoint,
            "model": model,
            "original_input": f"Uploaded file: {uploaded_files[0]}"
        }
        
        return jsonify({
            "session_id": session_id,
            "structure": structure,
            "thumbnails": thumbnails,
            "download_url": f"/api/ppt/download/{session_id}"
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

